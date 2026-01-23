import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Design Constants ---
DARK_BLUE = RGBColor(11, 56, 96)   # #0B3860
LIGHT_BLUE = RGBColor(77, 166, 255) # #4DA6FF
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(128, 128, 128)
LIGHT_GREY = RGBColor(240, 240, 240)
ACCENT_GREEN = RGBColor(46, 204, 113) # #2ECC71

def set_background(slide, color):
    # Set solid background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, DARK_BLUE)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Aceleração de QA com Inteligência e Automação"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle.text = "Extensão BDD v1.4\nAnálise de Impacto e ROI"
    subtitle.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

def add_content_slide(prs, title_text, bullet_points, highlight_box=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, WHITE)
    
    # Custom Title Style
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    title.top = Inches(0.5)
    
    # Content body
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Pt(14)
        
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)

    # Optional Highlight Box
    if highlight_box:
        left = Inches(5.5)
        top = Inches(2.5)
        width = Inches(4.0)
        height = Inches(3.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.color.rgb = LIGHT_BLUE
        
        text_frame = shape.text_frame
        text_frame.text = highlight_box
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.color.rgb = WHITE
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.word_wrap = True

def add_comparison_table_slide(prs):
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, WHITE)
    
    title = slide.shapes.title
    title.text = "Comparativo: Manual vs Tradicional vs Sistema BDD"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    # Table data
    rows = 4
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(4.0)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Critério", "Teste Manual", "Automação Tradicional", "Sistema BDD (Extensão)"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data
    data = [
        ["Setup / Criação", "Zero (imediato)", "Alto (4-8h/cenário)", "Mínimo (15min/cenário)"],
        ["Esforço Cognitivo", "Alto (repetitivo)", "Alto (código complexo)", "Baixo (gravação)"],
        ["Manutenção", "N/A (Refazer tudo)", "Média/Alta", "Baixa (POM Gerado)"]
    ]
    
    for r, row_data in enumerate(data, start=1):
        for c, value in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if c == 3: # Highlight the system column
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 240, 255)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

def add_roi_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, DARK_BLUE)
    
    title = slide.shapes.title
    title.text = "Impacto Financeiro e ROI"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.bold = True
    
    # Left content text
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Estimativa para 100 Cenários Automatizados:"
    tf.paragraphs[0].font.color.rgb = LIGHT_BLUE
    
    p = tf.add_paragraph()
    p.text = "• Método Tradicional: ~400 Horas"
    p.font.color.rgb = WHITE
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "• Com Extensão BDD: ~25 Horas"
    p.font.color.rgb = ACCENT_GREEN
    p.font.bold = True
    p.font.size = Pt(24)

    # Big Impact Statement
    left = Inches(5)
    top = Inches(3)
    width = Inches(4.5)
    height = Inches(2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    
    text_frame = shape.text_frame
    text_frame.text = "ECONOMIA DE 93% NO TEMPO DE DESENVOLVIMENTO DE TESTES\n\nDe 2.5 Meses para < 1 Semana"
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.size = Pt(24)

def add_tech_gains_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, WHITE)
    
    title = slide.shapes.title
    title.text = "Ganhos Técnicos e Diferenciais"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    points = [
        "Governança de Arquitetura: Código gerado segue estritamente Page Object Model.",
        "Resiliência: Implementação automática de 'Wait' explícitos (zero flaky tests).",
        "Documentação Viva: Saída nativa em Gherkin (BDD), legível por todos.",
        "Sem Lock-in: Código Python/Selenium puro, propriedade da empresa."
    ]
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = points[0]
    for point in points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.space_before = Pt(20)

    for p in tf.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(50, 50, 50)

def main():
    prs = Presentation()
    
    # 1. Slide Título
    add_title_slide(prs)
    
    # 2. Desafio
    add_content_slide(prs, 
                      "O Desafio Atual", 
                      [
                          "Gargalos no QA Manual: Execução lenta e propensa a erros humanos.",
                          "Alto Custo de Automação Tradicional: Requer engenheiros seniores e longo tempo de setup.",
                          "Manutenção Complexa: Scripts manuais tendem a se tornar frágeis ('flaky') sem arquitetura rígida."
                      ],
                      highlight_box="Problema:\nComo escalar a qualidade sem aumentar exponencialmente os custos?")
    
    # 3. A Solução
    add_content_slide(prs,
                      "A Solução: Extensão BDD v1.4",
                      [
                          "Conceito: Transformar navegação manual diretamente em código profissional.",
                          "Automatiza a Automação: O QA foca no negócio, a ferramenta gera o código.",
                          "Output Profissional: Gherkin + Python + Selenium com Page Objects."
                      ],
                      highlight_box="Inovação:\nReduz a barreira técnica. Um QA Manual entrega automação no Dia 1.")

    # 4. Comparativo Tabela
    add_comparison_table_slide(prs)

    # 5. Tempo e Velocidade
    add_content_slide(prs,
                      "Velocidade: Time-to-Market",
                      [
                          "Automação Tradicional: 4 Horas por Cenário",
                          "- Análise, Inspeção de Elementos, Page Objects, Código.",
                          "",
                          "Com Extensão BDD: 15 Minutos por Cenário",
                          "- Gravação (5min) + Exportação (5min) + Validação (5min)."
                      ],
                      highlight_box="Aceleração de 16x\nna entrega de novos testes.")
    
    # 6. ROI
    add_roi_slide(prs)
    
    # 7. Ganhos Técnicos
    add_tech_gains_slide(prs)

    # 8. DevOps & CI/CD (New)
    add_content_slide(prs,
                      "Integração DevOps & CI/CD",
                      [
                          "Esteira Automática: O código gerado é compatível nativamente com Jenkins, GitHub Actions, Azure DevOps.",
                          "Execução Headless: Testes prontos para rodar em containers Docker sem interface gráfica.",
                          "Feedback Rápido: Integração contínua garante que falhas sejam detectadas minutos após o commit."
                      ],
                      highlight_box="Plug & Play:\nGere o teste localmente, execute na nuvem imediatamente.")

    # 9. Benefícios Intangíveis (New)
    add_content_slide(prs,
                      "Cultura, Retenção e Propriedade",
                      [
                          "Retenção de Talentos: QAs deixam de ser 'robôs de teste' e viram Engenheiros de Qualidade analíticos.",
                          "Propriedade Intelectual (No Vendor Lock-in):",
                          "- O código Python/Selenium gerado é 100% ativo da empresa.",
                          "- Sem licenças mensais caras ou dependência de plataformas proprietárias."
                      ],
                      highlight_box="Segurança Jurídica:\nSeu ativo de automação é seu, para sempre.")
    
    # 10. Conclusão
    slide_layout = prs.slide_layouts[1] # Title Only like
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    set_background(slide, DARK_BLUE)
    
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Transforme seu QA de um gargalo operacional\npara um parceiro estratégico de velocidade."
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(32)

    output_file = "Apresentacao_Executiva_BDD_System.pptx"
    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")

if __name__ == "__main__":
    main()
